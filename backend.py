"""
DART Analytics Backend Server
============================

A Flask-based web server that provides REST API endpoints for generating 
MSD (Moving Standard Deviation) control charts from CSV/Excel data files.

Key Features:
- File upload and parsing (CSV/Excel)
- Automatic column type detection
- Control chart generation with statistical analysis
- PowerPoint export functionality
- Session management for multi-chart workflows
- Generates Excel file with outlier cells highlighted.
"""

# Core Flask imports for web server functionality
from flask import Flask, request, jsonify, send_from_directory, send_file
from flask_cors import CORS  # Enable cross-origin requests from frontend

# Data processing and analysis libraries
import pandas as pd  # Data manipulation and analysis
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend for server environment
import matplotlib.pyplot as plt  # Chart generation
import numpy as np  # Numerical computations

# System and utility imports
import os  # File system operations
import io  # Input/output operations for in-memory files
import base64  # Encoding images for JSON transmission
from datetime import datetime, timedelta  # Date/time handling
import traceback  # Error debugging
import warnings  # Warning management
import re  # Regular expressions for data cleaning
import csv  # CSV file parsing
import chardet  # Character encoding detection

# PowerPoint generation libraries
from pptx import Presentation  # Create PowerPoint presentations
from pptx.util import Inches, Pt  # Measurement units
from pptx.enum.text import PP_ALIGN  # Text alignment
from pptx.dml.color import RGBColor  # Color management

# Excel styling library
import openpyxl
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
from openpyxl.utils import get_column_letter

# Performance monitoring
from performance_monitor import monitor

# Suppress matplotlib warnings in production
warnings.filterwarnings("ignore", "Could not infer format", UserWarning)

# Initialize Flask application
app = Flask(__name__, static_folder='.', static_url_path='')
CORS(app)  # Enable CORS for all routes to allow frontend communication

def convert_numpy_types(obj):
    """
    Recursively convert NumPy data types to native Python types for JSON serialization.
    
    This function is essential because NumPy types (int64, float64, bool_, etc.) 
    are not JSON serializable by default. It handles nested structures like 
    dictionaries and lists containing NumPy types.
    
    Args:
        obj: Any object that may contain NumPy types
        
    Returns:
        Object with all NumPy types converted to native Python types
        
    Example:
        >>> convert_numpy_types(np.int64(42))
        42
        >>> convert_numpy_types({'mean': np.float64(3.14)})
        {'mean': 3.14}
    """
    if isinstance(obj, np.integer): 
        return int(obj)
    if isinstance(obj, np.floating): 
        return float(obj)
    if isinstance(obj, np.bool_): 
        return bool(obj)
    if isinstance(obj, np.ndarray): 
        return obj.tolist()
    if isinstance(obj, dict): 
        return {key: convert_numpy_types(value) for key, value in obj.items()}
    if isinstance(obj, list): 
        return [convert_numpy_types(item) for item in obj]
    return obj

class CheckColumns:
    """
    Main class for generating MSD (Moving Standard Deviation) control charts.
    
    This class handles the core analytics functionality including:
    - Data preprocessing and validation
    - Statistical calculations for control limits
    - Chart generation with matplotlib
    - Outlier detection and analysis
    
    The class implements statistical process control principles to identify
    data points that fall outside normal variation patterns.
    """
    
    def __init__(self, data_df):
        if not isinstance(data_df, pd.DataFrame):
            raise ValueError("Input data must be a pandas DataFrame.")
        self.df = data_df.copy()
        self.session_charts = []

    def create_control_chart(self, value_col, date_col, cut_cols=None, filters=None, aggregation_period='W', rolling_window=7, std_dev=2):
        """
        Generate MSD control charts from the loaded data.
        
        This is the main method that orchestrates the entire chart generation process:
        1. Applies user-specified filters to the data
        2. Cleans and validates the data (removes nulls, converts types)
        3. Groups data by cut columns if specified
        4. Generates individual charts for each group
        5. Returns chart data and processing statistics
        
        Args:
            value_col (str): Column name containing the numeric values to analyze
            date_col (str or list): Column name(s) containing date/time information
                                    Can be a single column or list of columns to combine
            cut_cols (list, optional): Columns to group data by (creates separate charts)
            filters (dict, optional): Column filters to apply before analysis
            aggregation_period (str): Time period for data aggregation ('D', 'W', 'M', 'Y')
            rolling_window (int): Number of periods for rolling average calculation
            std_dev (float): Number of standard deviations for control limits
            
        Returns:
            tuple: (charts_data, status_report)
                - charts_data: List of chart dictionaries with images and statistics
                - status_report: Dictionary with processing information
        """
        # Clear any existing matplotlib state to prevent memory leaks
        plt.close('all')
        plt.clf()
        plt.cla()
        
        # Step 1: Apply user-specified filters
        filtered_df = self.df
        if filters:
            for col, values in filters.items():
                if col in filtered_df.columns and values:
                    filtered_df = filtered_df[filtered_df[col].isin(values)]
        
        if filtered_df.empty: 
            return [], {"status": "No data remaining after applying filters."}

        # Step 1.5: Handle multiple date columns (combine them)
        combined_date_col = '_combined_timeseries_'
        original_date_cols = []
        
        if isinstance(date_col, list) and len(date_col) > 1:
            # Multiple columns selected - need to combine them
            original_date_cols = date_col.copy()
            
            # Verify all columns exist
            missing_cols = [col for col in date_col if col not in filtered_df.columns]
            if missing_cols:
                return [], {"status": f"Error: Columns not found: {', '.join(missing_cols)}"}
            
            # Try to combine the columns
            try:
                # Create a copy to work with
                temp_df = filtered_df[date_col].copy()
                
                # Check if columns can be combined into a datetime
                # Strategy 1: Try concatenating as strings and parsing
                combined_str = temp_df.astype(str).agg(' '.join, axis=1)
                combined_datetime = pd.to_datetime(combined_str, errors='coerce')
                
                # Check success rate
                success_rate = combined_datetime.notna().sum() / len(combined_datetime)
                
                if success_rate > 0.7:  # At least 70% successfully combined
                    filtered_df[combined_date_col] = combined_datetime
                    date_col = combined_date_col
                else:
                    # Strategy 2: Try numeric combination (e.g., Year + Month)
                    all_numeric = all(pd.api.types.is_numeric_dtype(temp_df[col]) for col in date_col)
                    
                    if all_numeric:
                        # Combine as concatenated integer (e.g., Year=2023, Month=12 -> 202312)
                        combined_numeric = temp_df.astype(str).agg(''.join, axis=1).astype(float)
                        filtered_df[combined_date_col] = combined_numeric
                        date_col = combined_date_col
                    else:
                        return [], {
                            "status": f"Error: Cannot combine columns {', '.join(original_date_cols)}. "
                                     f"Only {success_rate*100:.1f}% of values could be combined into valid dates. "
                                     f"Please ensure columns contain compatible date/time components (e.g., Year + Month, or Date + Time)."
                        }
                        
            except Exception as e:
                return [], {
                    "status": f"Error combining date columns {', '.join(original_date_cols)}: {str(e)}. "
                             f"Please select a single time series column or compatible columns that can be combined."
                }
        elif isinstance(date_col, list) and len(date_col) == 1:
            date_col = date_col[0]  # Single column in list format

        # Step 2: Data cleaning and validation
        initial_rows = len(filtered_df)
        clean_df = filtered_df.dropna(subset=[date_col, value_col]).copy()
        
        clean_df[value_col] = pd.to_numeric(clean_df[value_col], errors='coerce')
        
        # Check if date column is numeric (like week numbers: 12301, 12302, etc.)
        # If numeric, keep as-is for sorting; otherwise convert to datetime
        if pd.api.types.is_numeric_dtype(clean_df[date_col]):
            # Keep numeric time series column as-is (e.g., week numbers)
            clean_df[date_col] = pd.to_numeric(clean_df[date_col], errors='coerce')
        else:
            # Try to convert to datetime for actual date columns
            clean_df[date_col] = pd.to_datetime(clean_df[date_col], errors='coerce')
        
        clean_df.dropna(subset=[date_col, value_col], inplace=True)
        
        # Sort by time series column to ensure chronological order
        clean_df = clean_df.sort_values(by=date_col).reset_index(drop=True)
        
        rows_dropped = initial_rows - len(clean_df)
        if clean_df.empty: 
            return [], {"status": f"{rows_dropped} rows dropped; no data remains."}

        charts_data = []
        
        # Step 3: Handle data grouping (cut columns)
        if cut_cols and len(cut_cols) > 0:
            valid_cut_cols = []
            for col in cut_cols:
                if col in clean_df.columns and clean_df[col].nunique() > 1:
                    valid_cut_cols.append(col)
            
            if valid_cut_cols:
                grouped_by_cuts = clean_df.groupby(valid_cut_cols)
                
                for group_key, group_df in grouped_by_cuts:
                    if len(group_df) < 2:
                        continue
                    
                    chart_data = self._process_group(group_df, value_col, date_col, group_key, valid_cut_cols, aggregation_period, rolling_window, std_dev)
                    if chart_data:
                        charts_data.append(chart_data)
            else:
                chart_data = self._process_group(clean_df, value_col, date_col, None, [], aggregation_period, rolling_window, std_dev)
                if chart_data:
                    charts_data.append(chart_data)
        else:
            chart_data = self._process_group(clean_df, value_col, date_col, None, [], aggregation_period, rolling_window, std_dev)
            if chart_data:
                charts_data.append(chart_data)

        return charts_data, {"status": f"Analysis complete. {rows_dropped} rows excluded. Generated {len(charts_data)} charts."}

    def _process_group(self, group_df, value_col, date_col, group_key, cut_cols, aggregation_period, rolling_window, std_dev):
        """
        Process a single data group and calculate statistical control parameters.
        
        This method performs the core statistical analysis for MSD control charts:
        1. Aggregates data by time period (daily, weekly, monthly, yearly)
        2. Calculates sigma estimate using moving range method
        3. Computes rolling mean and control limits
        4. Identifies outliers and zero values
        5. Generates the visual chart
        
        Args:
            group_df (pd.DataFrame): Data for this specific group
            value_col (str): Column containing values to analyze
            date_col (str): Column containing date information
            group_key: Identifier for this group (from cut columns)
            cut_cols (list): List of grouping column names
            aggregation_period (str): Time aggregation period
            rolling_window (int): Window size for rolling calculations
            std_dev (float): Standard deviation multiplier for control limits
            
        Returns:
            dict or None: Chart data dictionary or None if processing fails
        """
        try:
            # Store original indices for Excel highlighting
            original_indices = group_df.index.tolist()
            
            # Sort data by time series column before processing
            group_df = group_df.sort_values(by=date_col).reset_index(drop=True)
            
            # Check if date column is numeric (like week numbers) - cannot resample numeric columns
            is_numeric_timeseries = pd.api.types.is_numeric_dtype(group_df[date_col]) and not pd.api.types.is_datetime64_any_dtype(group_df[date_col])
            
            # Aggregate data by time period (or use original data if NONE or numeric time series)
            if aggregation_period == 'NONE' or is_numeric_timeseries:
                # Use original data without aggregation
                resampled_df = group_df[[date_col, value_col]].copy()
                resampled_df = resampled_df.sort_values(by=date_col).reset_index(drop=True)
            else:
                # Sort before resampling to ensure chronological order (datetime columns only)
                group_df_sorted = group_df.sort_values(by=date_col)
                group_df_resampled = group_df_sorted.set_index(date_col)[[value_col]].resample(aggregation_period).mean().dropna().reset_index()
                if len(group_df_resampled) < 2:
                    return None
                resampled_df = group_df_resampled.copy()
            
            if len(resampled_df) < 2:
                return None
            
            # Calculate sigma estimate
            non_zero_values = resampled_df[value_col].replace(0, np.nan)
            sigma_estimate = np.nan
            if len(non_zero_values.dropna()) > 1:
                moving_ranges = non_zero_values.dropna().diff().abs()
                sigma_estimate = (moving_ranges.mean() / 1.128) if len(moving_ranges.dropna()) > 1 else non_zero_values.std(ddof=1)
            
            if pd.isna(sigma_estimate) or sigma_estimate == 0: 
                return None

            # Calculate control parameters
            rolling_mean = non_zero_values.rolling(window=rolling_window, min_periods=1).mean()
            resampled_df['rolling_mean'] = rolling_mean.interpolate(method='linear', limit_direction='both')
            resampled_df['upper_bound'] = resampled_df['rolling_mean'] + (std_dev * sigma_estimate)
            resampled_df['lower_bound'] = (resampled_df['rolling_mean'] - (std_dev * sigma_estimate)).clip(lower=0)
            
            # Enhanced outlier detection
            resampled_df['high_outlier'] = resampled_df.apply(
                lambda r: r[value_col] if r[value_col] > r['upper_bound'] and r[value_col] != 0 else np.nan,
                axis=1
            )
            resampled_df['low_outlier'] = resampled_df.apply(
                lambda r: r[value_col] if r[value_col] < r['lower_bound'] and r[value_col] != 0 else np.nan,
                axis=1
            )
            
            # Extreme outliers
            resampled_df['extreme_high_outlier'] = resampled_df.apply(
                lambda r: r[value_col] if (r[value_col] > r['rolling_mean'] + (5 * (r['upper_bound'] - r['rolling_mean']))) else np.nan,
                axis=1
            )
            resampled_df['extreme_low_outlier'] = resampled_df.apply(
                lambda r: r[value_col] if (r[value_col] < r['rolling_mean'] - (5 * (r['rolling_mean'] - r['lower_bound'])) 
                                         and r[value_col] != 0) else np.nan,
                axis=1
            )
            
            # Combined outliers
            resampled_df['outlier'] = resampled_df.apply(
                lambda r: r[value_col] if pd.notna(r['high_outlier']) or pd.notna(r['low_outlier']) else np.nan,
                axis=1
            )
            
            # Zero values
            resampled_df['zero_value'] = resampled_df.apply(
                lambda r: r[value_col] if r[value_col] == 0 else np.nan,
                axis=1
            )
            
            # Generate chart
            chart_data = self.generate_plot(resampled_df.copy(), value_col, date_col, group_key, cut_cols, rolling_window, std_dev, aggregation_period)
            
            # Add outlier mapping information and aggregation period
            if chart_data:
                chart_data['resampled_data'] = resampled_df.to_dict('records')
                chart_data['aggregation_period'] = aggregation_period
            
            plt.close('all')
            plt.clf()
            plt.cla()
            
            return chart_data
            
        except Exception as e:
            print(f"Error processing group {group_key}: {e}")
            return None

    def generate_plot(self, plot_df, value_col, date_col, group_key, cut_cols, rolling_window, std_dev, aggregation_period):
        """
        Generate a professional MSD control chart using matplotlib.
        
        This method creates a comprehensive statistical process control chart with:
        - Actual data points connected by lines
        - Rolling mean (center line) 
        - Upper and lower control limits
        - Highlighted outliers and zero values
        - Professional styling and annotations
        
        Args:
            plot_df (pd.DataFrame): Processed data ready for plotting
            value_col (str): Name of the value column being analyzed
            date_col (str): Name of the date column
            group_key: Identifier for data grouping
            cut_cols (list): Column names used for grouping
            rolling_window (int): Rolling window size for display
            std_dev (float): Standard deviation multiplier for display
            aggregation_period (str): Time aggregation period for display
            
        Returns:
            dict: Chart data including base64 image, statistics, and metadata
        """
        # Step 1: Initialize clean matplotlib environment
        plt.clf()
        plt.cla()
        plt.close('all')
        
        plt.style.use('seaborn-v0_8-whitegrid')
        fig, ax = plt.subplots(figsize=(15, 7))
        ax.clear()
        
        # Calculate dynamic y-axis limits
        all_values = pd.concat([
            plot_df[value_col],
            plot_df['upper_bound'],
            plot_df['lower_bound'].replace(0, np.nan)
        ]).dropna()
        
        if len(all_values) > 0:
            y_min = all_values.min()
            y_max = all_values.max()
            y_range = y_max - y_min
            y_padding = y_range * 0.1
            y_min = max(0, y_min - y_padding)
            y_max = y_max + y_padding
            ax.set_ylim(y_min, y_max)
        
        # Plot data
        ax.plot(plot_df[date_col], plot_df[value_col], 
               marker='o', linestyle='-', color='blue', 
               label='Actual', markersize=4)
        
        ax.plot(plot_df[date_col], plot_df['rolling_mean'], 
               color='green', linestyle='-', linewidth=2, 
               label=f'{rolling_window}-period Rolling Mean')
        
        ax.plot(plot_df[date_col], plot_df['upper_bound'], 
               color='red', linestyle='--', linewidth=1, 
               label=f'±{std_dev}σ Control Limit')
        ax.plot(plot_df[date_col], plot_df['lower_bound'], 
               color='red', linestyle='--', linewidth=1)
        
        ax.fill_between(plot_df[date_col], plot_df['lower_bound'], plot_df['upper_bound'], 
                       color='gray', alpha=0.15, label='Control Zone')
        
        # Highlight outliers
        high_outliers_mask = plot_df['high_outlier'].notna()
        low_outliers_mask = plot_df['low_outlier'].notna()
        extreme_high_mask = plot_df['extreme_high_outlier'].notna()
        extreme_low_mask = plot_df['extreme_low_outlier'].notna()
        zero_mask = plot_df['zero_value'].notna()
        
        regular_outliers_mask = ((high_outliers_mask | low_outliers_mask) & 
                               ~(extreme_high_mask | extreme_low_mask))
        if regular_outliers_mask.any():
            regular_values = plot_df.loc[regular_outliers_mask, 'outlier']
            ax.scatter(plot_df.loc[regular_outliers_mask, date_col],
                      regular_values,
                      color='red', s=80, zorder=5, alpha=0.7,
                      label=f'Outliers ({regular_outliers_mask.sum()})')
        
        if extreme_high_mask.any():
            ax.scatter(plot_df.loc[extreme_high_mask, date_col],
                      plot_df.loc[extreme_high_mask, 'extreme_high_outlier'],
                      color='red', s=200, marker='*', zorder=7, 
                      edgecolors='darkred', linewidth=1,
                      label=f'Severe High Outliers ({extreme_high_mask.sum()})')
            
        if extreme_low_mask.any():
            ax.scatter(plot_df.loc[extreme_low_mask, date_col],
                      plot_df.loc[extreme_low_mask, 'extreme_low_outlier'],
                      color='red', s=200, marker='*', zorder=7,
                      edgecolors='darkred', linewidth=1,
                      label=f'Severe Low Outliers ({extreme_low_mask.sum()})')
        
        if zero_mask.any():
            ax.scatter(plot_df.loc[zero_mask, date_col],
                      plot_df.loc[zero_mask, 'zero_value'],
                      color='orange', s=50, zorder=5,
                      label=f'Zero Values ({zero_mask.sum()})')
        
        # Create title
        title_add, group_name = "", "All Data"
        if cut_cols and group_key is not None:
            if isinstance(group_key, tuple):
                group_details = list(zip(cut_cols, group_key))
            else:
                group_details = [(cut_cols[0], group_key)]
            title_add = " - " + ", ".join(f"{col}={val}" for col, val in group_details)
            group_name = title_add.replace(" - ", "")
        
        agg_map = {'D': 'Daily', 'W': 'Weekly', 'M': 'Monthly', 'Y': 'Yearly', 'NONE': 'Original Data'}
        
        # Create title based on aggregation type
        if aggregation_period == 'NONE':
            title_text = f"MSD Control Chart for {value_col}{title_add}\n(Original Data, {rolling_window}-period rolling window)"
            xlabel = 'Time (Original Data Points)'
        else:
            title_text = f"MSD Control Chart for {value_col}{title_add}\n({agg_map.get(aggregation_period, '')} Aggregation, {rolling_window}-period rolling window)"
            xlabel = f'Time ({agg_map.get(aggregation_period, "")})'
        
        ax.set_title(title_text, fontsize=14, weight='bold')
        ax.set_xlabel(xlabel, fontsize=12)
        ax.set_ylabel(value_col, fontsize=12)
        
        plt.xticks(rotation=30, ha="right")
        plt.tight_layout()
        ax.legend(loc='best')
        
        # Save chart
        img_buffer = io.BytesIO()
        plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
        
        plt.close(fig)
        plt.clf()
        plt.cla()
        
        # Calculate statistics
        total_points = len(plot_df)
        half_point = max(1, total_points // 2)
        
        total_outliers = int((high_outliers_mask | low_outliers_mask).sum())
        total_extreme_outliers = int((extreme_high_mask | extreme_low_mask).sum())
        recent_outliers = int((high_outliers_mask | low_outliers_mask).iloc[half_point:].sum())
        recent_extreme_outliers = int((extreme_high_mask | extreme_low_mask).iloc[half_point:].sum())
        
        y_min, y_max = ax.get_ylim()
        off_scale_high = int(sum(plot_df[value_col] > y_max))
        off_scale_low = int(sum((plot_df[value_col] < y_min) & (plot_df[value_col] != 0)))
        
        all_outliers = {
            'total': total_outliers,
            'high': int(high_outliers_mask.sum()),
            'low': int(low_outliers_mask.sum()),
            'extreme_high': int(extreme_high_mask.sum()),
            'extreme_low': int(extreme_low_mask.sum()),
            'recent_total': recent_outliers,
            'recent_extreme': recent_extreme_outliers,
            'max_value': float(plot_df[value_col].max()),
            'min_value': float(plot_df[value_col].min()),
            'visible_range': {
                'min': float(y_min),
                'max': float(y_max)
            },
            'off_scale': {
                'high': off_scale_high,
                'low': off_scale_low,
                'total': off_scale_high + off_scale_low
            }
        }
        
        def safe_float(val):
            try:
                f = float(val)
                if pd.isna(f):
                    return None
                return f
            except Exception:
                return None

        chart_data = {
            'image': base64.b64encode(img_buffer.getvalue()).decode(),
            'title': f"Chart for {value_col}{title_add}", 
            'group': group_name, 
            'data_points': total_points,
            'outlier_stats': all_outliers,
            'outliers': total_outliers,
            'latter_half_outliers': recent_outliers,
            'zero_values': int(zero_mask.sum()),
            'statistics': {
                'mean': safe_float(plot_df[value_col].mean()),
                'std': safe_float(plot_df[value_col].std()),
                'min': safe_float(plot_df[value_col].min()),
                'max': safe_float(plot_df[value_col].max()),
                'control_limits': {
                    'upper': safe_float(plot_df['upper_bound'].mean()),
                    'lower': safe_float(plot_df['lower_bound'].replace(0, np.nan).mean())
                }
            }
        }

        chart_data['_image_buffer'] = img_buffer.getvalue()

        return chart_data

# ============================================================================
# GLOBAL SESSION MANAGEMENT
# ============================================================================

session_charts_storage = {}

def cleanup_temp_files(max_age_hours=1):
    """Clean up temporary files older than specified age."""
    temp_dir = 'temp_data'
    if not os.path.isdir(temp_dir): 
        return
    
    current_time = datetime.now()
    
    for filename in os.listdir(temp_dir):
        file_path = os.path.join(temp_dir, filename)
        try:
            if os.path.isfile(file_path):
                file_age = current_time - datetime.fromtimestamp(os.path.getmtime(file_path))
                if file_age > timedelta(hours=max_age_hours):
                    os.remove(file_path)
                    print(f"Cleaned up old temp file: {filename}")
        except Exception as e: 
            print(f"Error cleaning up file {filename}: {e}")

def highlight_outliers_excel(df, value_col, date_col, charts_data, session_id, filters=None):
    """
    Generate Excel file with outlier cells highlighted in yellow.
    
    Args:
        df: Original DataFrame
        value_col: Column with values to analyze
        date_col: Column with dates
        charts_data: List of chart dictionaries with outlier information
        session_id: Session identifier for output file
        filters: Applied filters (optional)
    
    Returns:
        tuple: (Path to generated Excel file, filename only)
    """
    try:
        print(f"Starting Excel generation for session {session_id}")
        print(f"Date column(s): {date_col}, Type: {type(date_col)}")
        
        # Create copy of original data
        df_export = df.copy()
        
        # Apply same filters that were used for charts
        if filters:
            for col, values in filters.items():
                if col in df_export.columns and values:
                    df_export = df_export[df_export[col].isin(values)]
        
        # Handle date_col - it might be a list (multi-column) or string (single column)
        # Create a combined date column for matching
        is_multi_column = isinstance(date_col, list)
        
        if is_multi_column:
            # Multi-column time series - create combined column
            print(f"Multi-column date: {date_col}")
            combined_col_name = '_combined_date_'
            df_export[combined_col_name] = df_export[date_col].apply(
                lambda row: '_'.join(row.astype(str).values), axis=1
            )
            date_col_for_matching = combined_col_name
        else:
            # Single column - convert to datetime
            df_export[date_col] = pd.to_datetime(df_export[date_col], errors='coerce')
            date_col_for_matching = date_col
        
        df_export[value_col] = pd.to_numeric(df_export[value_col], errors='coerce')
        
        # Collect ALL outlier information from charts with proper aggregation period mapping
        all_outlier_indices = set()
        
        for chart in charts_data:
            if 'resampled_data' in chart:
                resampled_df = pd.DataFrame(chart['resampled_data'])
                
                # Debug: Print what columns are in resampled_df
                print(f"Resampled DataFrame columns: {resampled_df.columns.tolist()}")
                print(f"First row sample: {resampled_df.head(1).to_dict('records')}")
                
                # Get outlier rows from resampled data
                outlier_rows = resampled_df[resampled_df['outlier'].notna()]
                
                print(f"Chart '{chart.get('title', 'Unknown')}' has {len(outlier_rows)} outlier periods")
                
                for _, outlier_row in outlier_rows.iterrows():
                    # Get the date value from the outlier row
                    # The resampled data has a 'date' column that contains the combined date string
                    if 'date' in outlier_row:
                        outlier_date_value = outlier_row['date']
                    elif isinstance(date_col, list):
                        # Try to get from individual columns if they exist
                        if all(col in outlier_row.index for col in date_col):
                            outlier_date_value = '_'.join(str(outlier_row[col]) for col in date_col)
                        else:
                            print(f"Warning: Could not find date columns in outlier row. Available: {outlier_row.index.tolist()}")
                            continue
                    else:
                        # Single column: convert to datetime
                        outlier_date_value = pd.to_datetime(outlier_row[date_col])
                    
                    outlier_value = float(outlier_row[value_col])
                    
                    # For each outlier in the aggregated data, find ALL original rows
                    # that contributed to that aggregated period
                    
                    # Determine the aggregation period window
                    aggregation_period = chart.get('aggregation_period', 'W')
                    
                    if aggregation_period == 'NONE':
                        # Original data - exact match
                        if isinstance(date_col, list):
                            # Multi-column: match by combined string
                            matching_indices = df_export[
                                (df_export[date_col_for_matching] == outlier_date_value) &
                                (df_export[value_col] == outlier_value)
                            ].index.tolist()
                        else:
                            # Single column: match by datetime
                            matching_indices = df_export[
                                (df_export[date_col] == outlier_date_value) &
                                (df_export[value_col] == outlier_value)
                            ].index.tolist()
                        
                    elif aggregation_period == 'W':
                        # Weekly - find all rows in the same week
                        if is_multi_column:
                            # For multi-column, match by exact combined value
                            matching_indices = df_export[
                                df_export[date_col_for_matching] == outlier_date_value
                            ].index.tolist()
                        else:
                            week_start = outlier_date_value - pd.Timedelta(days=outlier_date_value.dayofweek)
                            week_end = week_start + pd.Timedelta(days=6)
                            
                            matching_indices = df_export[
                                (df_export[date_col_for_matching] >= week_start) &
                                (df_export[date_col_for_matching] <= week_end)
                            ].index.tolist()
                        
                    elif aggregation_period == 'D':
                        # Daily - find rows on the same day
                        if is_multi_column:
                            # For multi-column, match by exact combined value
                            matching_indices = df_export[
                                df_export[date_col_for_matching] == outlier_date_value
                            ].index.tolist()
                        else:
                            matching_indices = df_export[
                                df_export[date_col_for_matching].dt.date == outlier_date_value.date()
                            ].index.tolist()
                        
                    elif aggregation_period == 'M':
                        # Monthly - find all rows in the same month
                        if is_multi_column:
                            # For multi-column, match by exact combined value
                            matching_indices = df_export[
                                df_export[date_col_for_matching] == outlier_date_value
                            ].index.tolist()
                        else:
                            matching_indices = df_export[
                                (df_export[date_col_for_matching].dt.year == outlier_date_value.year) &
                                (df_export[date_col_for_matching].dt.month == outlier_date_value.month)
                            ].index.tolist()
                        
                    elif aggregation_period == 'Y':
                        # Yearly - find all rows in the same year
                        if is_multi_column:
                            # For multi-column, match by exact combined value
                            matching_indices = df_export[
                                df_export[date_col_for_matching] == outlier_date_value
                            ].index.tolist()
                        else:
                            matching_indices = df_export[
                                df_export[date_col_for_matching].dt.year == outlier_date_value.year
                            ].index.tolist()
                    else:
                        # Fallback to a wider time window
                        if is_multi_column:
                            # For multi-column, match by exact combined value
                            matching_indices = df_export[
                                df_export[date_col_for_matching] == outlier_date_value
                            ].index.tolist()
                        else:
                            matching_indices = df_export[
                                (df_export[date_col_for_matching] >= outlier_date_value - pd.Timedelta(days=7)) &
                                (df_export[date_col_for_matching] <= outlier_date_value + pd.Timedelta(days=7))
                            ].index.tolist()
                    
                    all_outlier_indices.update(matching_indices)
                    # Format the outlier value for logging
                    outlier_display = str(outlier_date_value) if is_multi_column else outlier_date_value.date()
                    print(f"  Outlier period {outlier_display} matched {len(matching_indices)} original rows")
        
        print(f"Total outlier data points: {len(outlier_rows) if 'outlier_rows' in locals() else 'N/A'}")
        print(f"Total original rows to highlight: {len(all_outlier_indices)}")
        
        # Create output directory
        os.makedirs('temp_exports', exist_ok=True)
        
        # Generate unique filename with timestamp to avoid overwriting
        import time
        timestamp = int(time.time() * 1000)  # Milliseconds for uniqueness
        excel_filename = f'outliers_{session_id}_{timestamp}.xlsx'
        output_path = os.path.join('temp_exports', excel_filename)
        
        # Write DataFrame to Excel
        df_export.to_excel(output_path, index=False, engine='openpyxl')
        
        # Load workbook for styling
        wb = openpyxl.load_workbook(output_path)
        ws = wb.active
        
        # Define styles
        yellow_fill = PatternFill(start_color='FFFF00', end_color='FFFF00', fill_type='solid')
        header_fill = PatternFill(start_color='4472C4', end_color='4472C4', fill_type='solid')
        header_font = Font(bold=True, color='FFFFFF', size=11)
        border = Border(
            left=Side(style='thin', color='000000'),
            right=Side(style='thin', color='000000'),
            top=Side(style='thin', color='000000'),
            bottom=Side(style='thin', color='000000')
        )
        
        # Style header row
        for cell in ws[1]:
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal='center', vertical='center')
            cell.border = border
        
        # Highlight ALL outlier rows
        for idx in all_outlier_indices:
            excel_row = idx + 2  # +2 because Excel is 1-indexed and has header row
            
            # Highlight the entire row
            for col in range(1, ws.max_column + 1):
                cell = ws.cell(row=excel_row, column=col)
                cell.fill = yellow_fill
                cell.border = border
        
        print(f"Highlighted {len(all_outlier_indices)} rows in Excel")
        
        # Auto-adjust column widths
        for column in ws.columns:
            max_length = 0
            column_letter = get_column_letter(column[0].column)
            
            for cell in column:
                try:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                except:
                    pass
            
            adjusted_width = min(max_length + 2, 50)
            ws.column_dimensions[column_letter].width = adjusted_width
        
        # Freeze header row
        ws.freeze_panes = 'A2'
        
        # Count total outlier periods across all charts
        total_outlier_periods = 0
        for chart in charts_data:
            if 'resampled_data' in chart:
                resampled_df = pd.DataFrame(chart['resampled_data'])
                outlier_count = resampled_df['outlier'].notna().sum()
                total_outlier_periods += outlier_count
        
        # Add a summary sheet
        summary_ws = wb.create_sheet('Outlier Summary', 0)
        summary_ws['A1'] = 'Outlier Detection Summary'
        summary_ws['A1'].font = Font(bold=True, size=14)
        
        summary_ws['A3'] = 'Total Outlier Periods Found:'
        summary_ws['B3'] = total_outlier_periods
        summary_ws['A4'] = 'Rows Highlighted:'
        summary_ws['B4'] = len(all_outlier_indices)
        summary_ws['A5'] = 'Generation Date:'
        summary_ws['B5'] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        
        summary_ws['A7'] = 'Note: Outlier cells are highlighted in yellow on the Data sheet.'
        summary_ws['A8'] = 'Each outlier period (from aggregated data) may include multiple original data rows.'
        summary_ws['A7'].font = Font(italic=True)
        summary_ws['A8'].font = Font(italic=True, size=9)
        
        # Save workbook
        wb.save(output_path)
        print(f"Excel file saved successfully: {output_path}")
        
        return output_path, excel_filename
        
    except Exception as e:
        print(f"Error creating Excel with highlighted outliers: {e}")
        traceback.print_exc()
        return None, None

def create_powerpoint_export(charts_data, session_id):
    """Create a professional PowerPoint presentation from generated charts."""
    try:
        sorted_charts = sorted(charts_data, 
                             key=lambda x: x.get('latter_half_outliers', 0), 
                             reverse=True)
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        TITLE_FONT = 'Segoe UI'
        BODY_FONT = 'Segoe UI'
        BLUE = RGBColor(0, 102, 204)
        GRAY_BLUE_BG = RGBColor(240, 244, 248)
        DARK_TEXT = RGBColor(40, 40, 40)
        ACCENT = RGBColor(20, 80, 140)

        def style_textframe(frame, font_size=Pt(14), bold=False, color=DARK_TEXT):
            for p in frame.paragraphs:
                for run in p.runs:
                    run.font.name = BODY_FONT
                    run.font.size = font_size
                    run.font.bold = bold
                    run.font.color.rgb = color

        # Title Slide
        title_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = GRAY_BLUE_BG

        title_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.5), Inches(8), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.text = "DART Analytics Report"
        style_textframe(title_frame, Pt(40), True, ACCENT)

        subtitle_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.5), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = (f"Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}\n"
                               f"Total Charts: {len(sorted_charts)}")
        style_textframe(subtitle_frame, Pt(18), False, DARK_TEXT)

        # Executive Summary
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = GRAY_BLUE_BG

        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Executive Summary"
        style_textframe(title_frame, Pt(32), True, ACCENT)

        charts_with_recent_anomalies = len([c for c in sorted_charts if c.get('latter_half_outliers', 0) > 0])
        total_anomalies = sum(c.get('outliers', 0) for c in sorted_charts)
        recent_anomalies = sum(c.get('latter_half_outliers', 0) for c in sorted_charts)

        content_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(10.5), Inches(4))
        content_frame = content_box.text_frame
        content_frame.text = (
            f"Key Findings:\n\n"
            f"• Total Charts Analyzed: {len(sorted_charts)}\n"
            f"• Charts with Recent Anomalies: {charts_with_recent_anomalies}\n"
            f"• Total Anomalies Detected: {total_anomalies}\n"
            f"• Recent Anomalies (Last Half): {recent_anomalies}\n\n"
            f"Charts are ordered by recent anomaly count (highest first) for priority review."
        )
        style_textframe(content_frame, Pt(20), False, DARK_TEXT)

        # Individual Chart Slides
        for i, chart in enumerate(sorted_charts, 1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

            title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = f"{i}. {chart.get('title', 'Chart')}"
            style_textframe(title_frame, Pt(26), True, ACCENT)

            if '_image_buffer' in chart:
                img_stream = io.BytesIO(chart['_image_buffer'])
                slide.shapes.add_picture(img_stream, Inches(0.8), Inches(1.3), Inches(11.8), Inches(4.8))

            stats_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.8), Inches(1.2))
            stats_frame = stats_box.text_frame

            data_points = chart.get('data_points', 0)
            zero_pct = (chart.get('zero_values', 0) / data_points) * 100 if data_points > 0 else 0
            stats_text = (
                f"Statistics:\n"
                f"Data Points: {chart.get('data_points', 'N/A')} | "
                f"Total Outliers: {chart.get('outliers', 0)} | "
                f"Recent Outliers: {chart.get('latter_half_outliers', 0)} | "
                f"Zero Values: {chart.get('zero_values', 0)} ({zero_pct:.1f}%)\n"
                f"Mean: {chart.get('statistics', {}).get('mean', 0):.3f} | "
                f"Std Dev: {chart.get('statistics', {}).get('std', 0):.3f} | "
                f"Min: {chart.get('statistics', {}).get('min', 0):.3f} | "
                f"Max: {chart.get('statistics', {}).get('max', 0):.3f}\n\n"
                f"Note: Charts are prioritized by recent anomaly count."
            )
            stats_frame.text = stats_text
            style_textframe(stats_frame, Pt(14), False, DARK_TEXT)

            if chart.get('latter_half_outliers', 0) > 0:
                shape = slide.shapes.add_shape(
                    autoshape_type_id=1,
                    left=Inches(10.8),
                    top=Inches(0.3),
                    width=Inches(2.3),
                    height=Inches(0.6)
                )
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = ACCENT
                line = shape.line
                line.fill.background()

                text_frame = shape.text_frame
                text_frame.text = "High Priority Chart"
                style_textframe(text_frame, Pt(12), True, RGBColor(255, 255, 255))

        os.makedirs('temp_exports', exist_ok=True)
        ppt_filename = f"DART_Report_{session_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        ppt_path = os.path.join('temp_exports', ppt_filename)
        prs.save(ppt_path)

        print(f"PowerPoint report successfully created: {ppt_filename}")
        return ppt_path, ppt_filename
        
    except Exception as e:
        print(f"Error creating PowerPoint: {e}")
        traceback.print_exc()
        return None, None

# ============================================================================
# REST API ENDPOINTS
# ============================================================================

@app.route('/api/upload', methods=['POST'])
def upload_file():
    """
    Handle file upload and data preprocessing.
    
    This endpoint processes CSV and Excel files, performs data quality analysis,
    and prepares the data for chart generation. It includes:
    - Automatic encoding detection for CSV files
    - Column type inference (numeric, date, categorical)
    - Data quality reporting
    - Session management for multi-chart workflows
    
    Expected Request:
        POST /api/upload
        Content-Type: multipart/form-data
        Body: file (CSV or Excel file)
        
    Returns:
        JSON response with:
        - session_id: Unique identifier for this data session
        - filename: Original filename
        - rows/columns: Data dimensions
        - columns_info: Detailed column analysis
        - filter_options: Available filter values for categorical columns
        - quality_report: Data quality metrics
        
    Error Codes:
        400: Bad request (no file, unsupported format, parsing error)
        500: Internal server error
    """
    try:
        cleanup_temp_files()
        
        if 'file' not in request.files: 
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '': 
            return jsonify({'error': 'No file selected'}), 400
        
        filename = file.filename.lower()
        try:
            if filename.endswith('.csv'):
                raw_data = file.read()
                file.seek(0)
                
                detected_encoding = chardet.detect(raw_data[:10000])['encoding'] or 'utf-8'
                
                try:
                    dialect = csv.Sniffer().sniff(raw_data[:4096].decode(detected_encoding))
                    df = pd.read_csv(io.StringIO(raw_data.decode(detected_encoding)), 
                                   sep=dialect.delimiter, low_memory=False)
                except (csv.Error, UnicodeDecodeError):
                    df = pd.read_csv(io.StringIO(raw_data.decode(detected_encoding)), 
                                   low_memory=False)
                    
            elif filename.endswith(('.xlsx', '.xls')):
                sheet_name = request.form.get('sheet_name')
                
                excel_file = pd.ExcelFile(file, engine='openpyxl')
                sheet_names = excel_file.sheet_names
                
                if len(sheet_names) > 1 and not sheet_name:
                    return jsonify({
                        'multiple_sheets': True,
                        'sheet_names': sheet_names,
                        'message': 'Please select a sheet to analyze'
                    }), 200
                
                sheet_to_use = sheet_name if sheet_name else sheet_names[0]
                df = pd.read_excel(excel_file, sheet_name=sheet_to_use, engine='openpyxl')
            else:
                return jsonify({'error': 'Unsupported file type. Please use CSV or Excel.'}), 400
                
        except Exception as e:
            return jsonify({'error': f"Failed to parse file. Details: {str(e)}"}), 400
        
        if df.empty: 
            return jsonify({'error': 'File is empty or could not be read.'}), 400
        
        initial_rows = len(df)
        df.dropna(how='all', inplace=True)
        rows_dropped = initial_rows - len(df)
        
        columns_info, filter_options = [], {}
        sample_df = df.head(100).copy()
        
        for col in df.columns:
            dtype = df[col].dtype
            
            is_numeric = pd.api.types.is_numeric_dtype(dtype)
            
            # Check if column is datetime or can be converted to datetime
            is_date_like = (pd.api.types.is_datetime64_any_dtype(dtype) or 
                           (pd.to_datetime(sample_df[col], errors='coerce').notna().sum() / len(sample_df) > 0.5 
                            if len(sample_df) > 0 else False))
            
            # Also mark integer columns as date_like if they could be sequential time series
            # (e.g., week numbers: 12301, 12302, 12303, or 1, 2, 3, etc.)
            if is_numeric and not is_date_like and pd.api.types.is_integer_dtype(dtype):
                unique_count = df[col].nunique()
                # If column has reasonable number of unique values and appears sequential
                if 2 <= unique_count <= 1000:
                    is_date_like = True  # Allow integer columns as time series
            
            if (not is_numeric and not is_date_like and 
                df[col].nunique() < 50 and df[col].nunique() > 1):
                filter_options[col] = sorted(df[col].dropna().unique().tolist())
            
            columns_info.append({
                'name': col, 
                'type': str(dtype), 
                'is_numeric': is_numeric, 
                'is_date_like': is_date_like
            })
        
        session_id = re.sub(r'\W+', '', str(datetime.now().timestamp()))
        os.makedirs('temp_data', exist_ok=True)
        df.to_pickle(os.path.join('temp_data', f'{session_id}.pkl'))
        
        # Log upload metrics
        file_size = len(raw_data) if filename.endswith('.csv') else file.content_length or 0
        file_type = 'csv' if filename.endswith('.csv') else 'excel'
        monitor.log_upload(file_size, file_type)
        
        response_data = {
            'session_id': session_id, 
            'filename': file.filename, 
            'rows': len(df), 
            'columns': len(df.columns),
            'columns_info': columns_info, 
            'filter_options': filter_options,
            'quality_report': {'empty_rows_dropped': rows_dropped}
        }
        
        return jsonify(convert_numpy_types(response_data))
        
    except Exception as e:
        traceback.print_exc()
        return jsonify({'error': f'An unexpected error occurred: {str(e)}'}), 500

def analyze_root_causes(df, value_col, date_col, charts_data, filters=None):
    """
    Perform root cause analysis to identify factors correlated with outliers.
    
    Args:
        df: Original DataFrame
        value_col: Value column name
        date_col: Date column(s)
        charts_data: Generated charts with outlier information
        filters: Applied filters
        
    Returns:
        dict: Root cause analysis results
    """
    try:
        root_causes = {
            'correlations': [],
            'patterns': [],
            'insights': []
        }
        
        # Get numeric columns for correlation analysis
        numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
        if value_col in numeric_cols:
            numeric_cols.remove(value_col)
        
        # Calculate correlations with value column
        if numeric_cols:
            correlations = []
            for col in numeric_cols:
                try:
                    corr = df[value_col].corr(df[col])
                    if not np.isnan(corr) and abs(corr) > 0.3:  # Significant correlation threshold
                        correlations.append({
                            'factor': col,
                            'correlation': round(float(corr), 3),
                            'strength': 'Strong' if abs(corr) > 0.7 else 'Moderate'
                        })
                except:
                    pass
            
            # Sort by absolute correlation
            correlations.sort(key=lambda x: abs(x['correlation']), reverse=True)
            root_causes['correlations'] = correlations[:5]  # Top 5
        
        # Analyze patterns in outliers
        for chart in charts_data:
            if 'resampled_data' in chart:
                resampled_df = pd.DataFrame(chart['resampled_data'])
                outlier_rows = resampled_df[resampled_df['outlier'].notna()]
                
                if len(outlier_rows) > 0:
                    # Check if outliers cluster in specific periods
                    if 'date' in outlier_rows.columns:
                        dates = outlier_rows['date'].tolist()
                        if len(dates) > 1:
                            root_causes['patterns'].append({
                                'type': 'temporal_clustering',
                                'description': f"Outliers detected in {len(dates)} periods",
                                'periods': dates[:5]  # First 5 periods
                            })
        
        # Generate insights
        if root_causes['correlations']:
            top_corr = root_causes['correlations'][0]
            root_causes['insights'].append({
                'type': 'correlation',
                'message': f"{top_corr['factor']} shows {top_corr['strength'].lower()} correlation ({top_corr['correlation']}) with {value_col}",
                'impact': 'high' if abs(top_corr['correlation']) > 0.7 else 'medium'
            })
        
        # Analyze outlier distribution
        total_outliers = sum(len(pd.DataFrame(chart.get('resampled_data', [])).query('outlier.notna()')) for chart in charts_data if 'resampled_data' in chart)
        extreme_outliers = sum(
            len(pd.DataFrame(chart.get('resampled_data', [])).query('extreme_high_outlier.notna() or extreme_low_outlier.notna()')) 
            for chart in charts_data if 'resampled_data' in chart
        )
        
        if total_outliers > 0:
            outlier_percentage = (extreme_outliers / total_outliers * 100) if total_outliers > 0 else 0
            if outlier_percentage > 20:
                root_causes['insights'].append({
                    'type': 'severity',
                    'message': f"{extreme_outliers} of {total_outliers} outliers ({outlier_percentage:.1f}%) are extreme deviations, indicating significant process variation",
                    'impact': 'high'
                })
            elif total_outliers > 10:
                root_causes['insights'].append({
                    'type': 'frequency',
                    'message': f"High outlier frequency detected ({total_outliers} total outliers). Consider investigating process stability",
                    'impact': 'medium'
                })
            else:
                root_causes['insights'].append({
                    'type': 'stability',
                    'message': f"Process shows relatively stable behavior with {total_outliers} outliers detected",
                    'impact': 'low'
                })
        
        # Check for categorical patterns
        categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
        if isinstance(date_col, list):
            categorical_cols = [c for c in categorical_cols if c not in date_col]
        elif isinstance(date_col, str):
            categorical_cols = [c for c in categorical_cols if c != date_col]
        
        # Analyze categorical factors if available
        if categorical_cols and total_outliers > 0:
            for col in categorical_cols[:2]:  # Top 2 categorical columns
                try:
                    category_counts = df[col].value_counts()
                    if len(category_counts) <= 10:  # Only if manageable number of categories
                        root_causes['insights'].append({
                            'type': 'categorical',
                            'message': f"Data contains {len(category_counts)} distinct {col} categories. Consider analyzing outliers by {col} segment",
                            'impact': 'medium'
                        })
                        break
                except:
                    pass
            
        return root_causes
        
    except Exception as e:
        print(f"Root cause analysis error: {e}")
        traceback.print_exc()
        return {'correlations': [], 'patterns': [], 'insights': []}

@app.route('/api/generate_chart', methods=['POST'])
def generate_chart_api():
    """
    Generate MSD control charts based on user specifications.
    
    This endpoint is the core of the analytics functionality. It takes user
    parameters and generates statistical control charts with outlier detection.
    
    Expected Request:
        POST /api/generate_chart
        Content-Type: application/json
        Body: {
            "session_id": "unique_session_identifier",
            "value_column": "column_name_with_numeric_data",
            "date_column": "column_name_with_date_data" OR ["col1", "col2"], 
            "cut_columns": ["optional_grouping_columns"],
            "filters": {"column": ["filter_values"]},
            "aggregation_period": "W|D|M|Y",
            "rolling_window": 7,
            "std_dev": 2.0
        }
        Note: date_column can be a single column name (string) or array of column names
              to combine (e.g., ["Year", "Month"] or ["Date", "Time"])
        
    Returns:
        JSON response with:
        - success: Boolean indicating success
        - charts: Array of chart objects with images and statistics
        - message: Status message with processing details
        
    Error Codes:
        400: Bad request (missing parameters, invalid data)
        404: Session not found (expired or invalid)
        500: Internal server error
    """
    try:
        data = request.get_json()
        session_id = re.sub(r'\W+', '', data.get('session_id', ''))
        value_col, date_col = data.get('value_column'), data.get('date_column')
        
        if not all([session_id, value_col, date_col]): 
            return jsonify({'error': 'Missing required parameters'}), 400
        
        try:
            rolling_window = int(data.get('rolling_window', 7))
            std_dev = float(data.get('std_dev', 2))
        except (ValueError, TypeError):
            return jsonify({'error': 'Rolling window and std dev must be numbers.'}), 400
        
        df_path = os.path.join('temp_data', f'{session_id}.pkl')
        if not os.path.exists(df_path): 
            return jsonify({'error': 'Session expired or invalid.'}), 404
        
        df = pd.read_pickle(df_path)
        
        checker = CheckColumns(df)
        charts_data, report = checker.create_control_chart(
            value_col=value_col, 
            date_col=date_col, 
            cut_cols=data.get('cut_columns'), 
            filters=data.get('filters'),
            aggregation_period=data.get('aggregation_period', 'W'), 
            rolling_window=rolling_window, 
            std_dev=std_dev
        )
        
        if not charts_data: 
            return jsonify({'error': f"No valid data to generate charts. {report.get('status', '')}"}), 400
        
        # Store charts for PowerPoint export
        if session_id not in session_charts_storage:
            session_charts_storage[session_id] = []
        session_charts_storage[session_id].extend(charts_data)
        
        # Generate Excel file with highlighted outliers
        excel_path = None
        excel_filename = None
        try:
            excel_path, excel_filename = highlight_outliers_excel(
                df, 
                value_col, 
                date_col, 
                charts_data, 
                session_id,
                filters=data.get('filters')
            )
            print(f"Excel file generated: {excel_path}")
        except Exception as e:
            print(f"Excel generation failed: {e}")
            traceback.print_exc()
        
        # Prepare JSON response
        charts_for_json = []
        for chart in charts_data:
            chart_copy = chart.copy()
            if '_image_buffer' in chart_copy:
                del chart_copy['_image_buffer']
            if 'resampled_data' in chart_copy:
                del chart_copy['resampled_data']
            if 'outlier_dates' in chart_copy:
                del chart_copy['outlier_dates']
            if 'outlier_values' in chart_copy:
                del chart_copy['outlier_values']
            charts_for_json.append(chart_copy)
        
        # Log chart generation metrics
        total_outliers = sum(chart.get('outlier_count', 0) for chart in charts_data)
        monitor.log_chart_generation('MSD', total_outliers)
        
        # Perform root cause analysis
        # COMMENTED OUT - Feature in development
        # root_cause_analysis = analyze_root_causes(df, value_col, date_col, charts_data, data.get('filters'))
        
        response_data = {
            'success': True, 
            'charts': convert_numpy_types(charts_for_json), 
            'message': f"Generated {len(charts_data)} chart(s). {report.get('status', '')}"
            # 'root_cause_analysis': convert_numpy_types(root_cause_analysis)
        }
        
        # Add Excel download info if generated successfully
        if excel_path and excel_filename and os.path.exists(excel_path):
            response_data['excel_file'] = excel_filename
            response_data['excel_ready'] = True
        
        return jsonify(response_data)
        
    except Exception as e:
        traceback.print_exc()
        return jsonify({'error': f'An unexpected error occurred: {str(e)}'}), 500

@app.route('/api/download_excel/<excel_filename>', methods=['GET'])
def download_excel(excel_filename):
    """Download the Excel file with highlighted outliers."""
    try:
        # Sanitize the filename to prevent directory traversal
        excel_filename = os.path.basename(excel_filename)
        
        # Ensure it's a valid Excel filename pattern
        if not excel_filename.startswith('outliers_') or not excel_filename.endswith('.xlsx'):
            return jsonify({'error': 'Invalid Excel filename.'}), 400
        
        excel_path = os.path.join('temp_exports', excel_filename)
        
        if not os.path.exists(excel_path):
            return jsonify({'error': 'Excel file not found. Please generate charts first.'}), 404
        
        # Get custom filename from query parameters if provided
        custom_filename = request.args.get('filename', excel_filename)
        if not custom_filename.endswith('.xlsx'):
            custom_filename += '_outliers.xlsx'
        
        return send_file(
            excel_path,
            as_attachment=True,
            download_name=custom_filename,  # Use the custom filename for download
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )
        
    except Exception as e:
        traceback.print_exc()
        return jsonify({'error': f'Error downloading Excel: {str(e)}'}), 500

@app.route('/api/export_ppt/<session_id>', methods=['GET'])
def export_powerpoint(session_id):
    """Export all charts to PowerPoint."""
    try:
        session_id = re.sub(r'\W+', '', session_id)
        
        if session_id not in session_charts_storage or not session_charts_storage[session_id]:
            return jsonify({'error': 'No charts found for this session. Generate some charts first.'}), 404
        
        charts_data = session_charts_storage[session_id]
        
        ppt_path, ppt_filename = create_powerpoint_export(charts_data, session_id)
        
        if not ppt_path or not os.path.exists(ppt_path):
            return jsonify({'error': 'Failed to create PowerPoint presentation'}), 500
        
        # Log export metrics
        monitor.log_export('PowerPoint', len(charts_data))
        
        return send_file(
            ppt_path,
            as_attachment=True,
            download_name=ppt_filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        traceback.print_exc()
        return jsonify({'error': f'Error exporting PowerPoint: {str(e)}'}), 500

@app.route('/api/clear_session/<session_id>', methods=['DELETE'])
def clear_session_charts(session_id):
    """
    Clear stored charts for a specific session.
    
    This endpoint removes all chart data associated with a session to free
    memory and allow users to start fresh analyses.
    
    Expected Request:
        DELETE /api/clear_session/{session_id}
        
    Returns:
        JSON response confirming session clearance
    """
    try:
        session_id = re.sub(r'\W+', '', session_id)
        
        if session_id in session_charts_storage:
            del session_charts_storage[session_id]
            
        return jsonify({'success': True, 'message': 'Session charts cleared'})
    except Exception as e:
        return jsonify({'error': f'Error clearing session: {str(e)}'}), 500

@app.route('/api/health', methods=['GET'])
def health_check():
    """
    Health check endpoint for monitoring.
    
    Returns application health status and performance metrics.
    Useful for load balancers, monitoring systems, and DevOps.
    
    Returns:
        JSON with health status and metrics
    """
    try:
        health_status = monitor.get_health_status()
        return jsonify({
            'status': 'healthy',
            'timestamp': datetime.now().isoformat(),
            'metrics': health_status
        })
    except Exception as e:
        return jsonify({
            'status': 'unhealthy',
            'error': str(e),
            'timestamp': datetime.now().isoformat()
        }), 500

@app.route('/api/statistics', methods=['GET'])
def get_statistics():
    """
    Get comprehensive usage statistics and analytics.
    
    Returns detailed metrics about application usage, performance,
    and system health. Useful for dashboards and analytics.
    
    Returns:
        JSON with usage statistics
    """
    try:
        stats = monitor.get_statistics()
        return jsonify({
            'success': True,
            'statistics': stats,
            'timestamp': datetime.now().isoformat()
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/')
def serve_frontend(): 
    """Serve the main frontend HTML file."""
    return send_from_directory('.', 'index.html')

# ============================================================================
# APPLICATION STARTUP
# ============================================================================

if __name__ == '__main__':
    """
    Start the Flask development server.
    
    This creates necessary directories and starts the server in debug mode
    for development. In production, use a proper WSGI server like Gunicorn.
    """
    # Ensure temporary data directory exists
    os.makedirs('temp_data', exist_ok=True)
    os.makedirs('temp_exports', exist_ok=True)
    
    app.run(debug=True, host='0.0.0.0', port=5000)